#!/usr/bin/env python3
"""
export_deck.py — export an HTML deck to PDF and/or screenshot-mode PPTX.

Usage:
    python tools/export_deck.py decks/<deck-name>.html --pdf
    python tools/export_deck.py decks/<deck-name>.html --pptx
    python tools/export_deck.py decks/<deck-name>.html --pdf --pptx

Outputs land next to the HTML:
    decks/<deck-name>.pdf
    decks/<deck-name>.pptx

PDF:   uses deck-stage's @media print path — one page per slide.
PPTX:  takes a PNG of each slide at 1920×1080, embeds full-bleed in a
       16:9 PowerPoint deck via python-pptx. Pixel-perfect, NOT editable.

Editable PPTX (native PowerPoint text boxes and shapes) is intentionally
out of scope for the local pipeline — it requires custom OOXML generation
the design environment provides. Use the HTML as your editable artifact.
"""

import argparse
import sys
import tempfile
from pathlib import Path
from playwright.sync_api import sync_playwright


def export_pdf(html_path: Path, pdf_path: Path) -> None:
    """Print to PDF using Chrome headless. deck-stage's @media print
    rule lays each slide out as its own page at the design size, so the
    result is one slide per page with no extra setup."""
    with sync_playwright() as p:
        browser = p.chromium.launch(channel="chrome")  # uses local Chrome; switch to p.chromium.launch() for bundled
        page = browser.new_page()
        page.goto(html_path.absolute().as_uri(), wait_until="networkidle")
        page.wait_for_timeout(500)
        page.emulate_media(media="print")
        page.pdf(
            path=str(pdf_path),
            width="1920px",
            height="1080px",
            print_background=True,
            margin={"top": "0", "right": "0", "bottom": "0", "left": "0"},
            prefer_css_page_size=True,
        )
        browser.close()
    print(f"  ✓ PDF written → {pdf_path}", file=sys.stderr)


def export_pptx(html_path: Path, pptx_path: Path) -> None:
    """Screenshot each slide at 1920×1080, embed each as a full-bleed
    image in a 16:9 PPTX. Pixel-perfect, not editable."""
    from pptx import Presentation
    from pptx.util import Emu

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)
        screenshots: list[Path] = []

        with sync_playwright() as p:
            browser = p.chromium.launch(channel="chrome")  # uses local Chrome; switch to p.chromium.launch() for bundled
            context = browser.new_context(viewport={"width": 1920, "height": 1080})
            page = context.new_page()
            page.goto(html_path.absolute().as_uri(), wait_until="networkidle")

            # Disable deck-stage scaling so each section is rendered at
            # authored 1920×1080.
            page.evaluate("document.querySelector('deck-stage').setAttribute('noscale', '')")
            page.wait_for_timeout(400)

            total = page.evaluate("document.querySelector('deck-stage').length")
            print(f"  Capturing {total} slide(s) at 1920×1080…", file=sys.stderr)

            for i in range(total):
                page.evaluate(f"document.querySelector('deck-stage').goTo({i})")
                page.wait_for_timeout(350)
                shot = tmp_dir / f"slide-{i+1:03d}.png"
                page.locator("deck-stage section[data-deck-active]").screenshot(path=str(shot))
                screenshots.append(shot)

            browser.close()

        # Build the PPTX. 16:9 at 1920×1080 = 13.333" × 7.5" = 12192000 × 6858000 EMU
        prs = Presentation()
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)
        blank_layout = prs.slide_layouts[6]  # blank

        for shot in screenshots:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(shot),
                left=0, top=0,
                width=prs.slide_width, height=prs.slide_height,
            )

        prs.save(str(pptx_path))

    print(f"  ✓ PPTX written → {pptx_path} (screenshot mode, not editable)", file=sys.stderr)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("html", type=Path, help="Path to deck HTML")
    ap.add_argument("--pdf", action="store_true", help="Export PDF")
    ap.add_argument("--pptx", action="store_true", help="Export PPTX (screenshot mode)")
    args = ap.parse_args()

    if not (args.pdf or args.pptx):
        print("ERROR: pass --pdf, --pptx, or both.", file=sys.stderr)
        return 2

    if not args.html.exists():
        print(f"ERROR: {args.html} does not exist", file=sys.stderr)
        return 2

    stem = args.html.with_suffix("")
    print(f"Exporting {args.html} …", file=sys.stderr)

    if args.pdf:
        export_pdf(args.html, stem.with_suffix(".pdf"))
    if args.pptx:
        export_pptx(args.html, stem.with_suffix(".pptx"))

    return 0


if __name__ == "__main__":
    sys.exit(main())
